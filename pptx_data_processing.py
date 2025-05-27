from pptx import Presentation
from docx import Document
from docx.shared import Inches
from docx.shared import Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.text.paragraph import Paragraph
from docx.oxml.shared import qn as qn_shared

import glob
import os
import win32com.client
from PIL import Image, ImageDraw, ImageFont
import re

import pytesseract

# Set path to tesseract.exe
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

def export_slides_as_images(pptx_path, output_dir):
    pptx_path = os.path.abspath(pptx_path)
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    ppt_app = win32com.client.Dispatch("PowerPoint.Application")
    ppt_app.Visible = 1
    presentation = ppt_app.Presentations.Open(pptx_path, WithWindow=False)
    os.makedirs(output_dir, exist_ok=True)
    presentation.SaveAs(output_dir, 17)  # 17 = export as PNG
    presentation.Close()
    ppt_app.Quit()


def remove_logo_and_extract_heading(image_path, save_main_path=None, save_heading_path=None, top_px=116):
    image = Image.open(image_path)
    width, height = image.size

    # Crop top and bottom
    cropped_main = image.crop((0, top_px, width, height))
    heading_crop = image.crop((0, 0, int(width * 0.73), top_px))
    
    if not save_main_path:
        save_main_path = image_path

    if not save_heading_path:
        base, ext = os.path.splitext(image_path)
        save_heading_path = f"{base}_heading{ext}"

    cropped_main.save(save_main_path)
    heading_crop.save(save_heading_path)
        
    return save_main_path, save_heading_path

def insert_paragraph_after(paragraph, text="", style=None):
    new_p = OxmlElement("w:p")
    paragraph._p.addnext(new_p)
    new_para = Paragraph(new_p, paragraph._parent)
    if text:
        run = new_para.add_run(text)
    if style:
        new_para.style = style
    return new_para

def overlay_heading_on_image(image_path, save_path):
    # Load image
    img = Image.open(image_path).convert("RGB")

    # Extract text using OCR
    extracted_text = pytesseract.image_to_string(img).strip()

    # Initialize drawing context
    draw = ImageDraw.Draw(img)

    # Set font (fallback to default if custom font not available)
    try:
        font = ImageFont.truetype("arial.ttf", size=36)
    except:
        font = ImageFont.load_default()

    # Draw text at the top center
    text_width, text_height = draw.textsize(extracted_text, font=font)
    img_width, _ = img.size
    text_x = (img_width - text_width) // 2
    draw.text((text_x, 10), extracted_text, fill="black", font=font)

    # Save new image
    img.save(save_path)
    return save_path, extracted_text


def insert_slide_content_at_placeholder(template_path, output_path, placeholder, headings_dict, image_dict):
    doc = Document(template_path)
    for i, para in enumerate(doc.paragraphs):
        if placeholder in para.text:
            # Store ref to placeholder
            placeholder_para = para
            break
    else:
        raise ValueError(f"Placeholder '{placeholder}' not found in document.")

    # Remove placeholder text from that paragraph
    placeholder_para.text = placeholder_para.text.replace(placeholder, "").strip()

    # Insert slides after placeholder
    current_para = placeholder_para

    # Build the new content (text + images)
    for slide_key in sorted(headings_dict.keys(), key=lambda x: int(x.split('_')[-1])):
        heading_text = headings_dict[slide_key]
        image_path = image_dict.get(slide_key)

        # Heading
        current_para = insert_paragraph_after(current_para)
        run = current_para.add_run(heading_text)
        run.bold = True
        run.font.size = Pt(12)

        # Image
        if image_path:
            new_p = insert_paragraph_after(current_para)
            run = new_p.add_run()
            run.add_picture(image_path, width=Inches(6))
            new_p.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
            current_para = new_p

        # Page break
        current_para = insert_paragraph_after(current_para)
        run = current_para.add_run()
        run.add_break()  # page break

    doc.save(output_path)
    print(f"✅ Manual saved at: {output_path}")


def extract_DAP_text_and_images(pptx_path, output_dir):
    """
    Process and analyze all uploaded documents and extract text and images from a PowerPoint file.
    
    Args:
        pptx_path (str): Path to PowerPoint file
        output_dir (str): Directory to save outputs

    Returns:
        dict: Analysis results for each document
    """
    global slide_image_map, slide_headings_text
    
    
    # Load PowerPoint presentation
    prs = Presentation(pptx_path)
    os.makedirs(output_dir, exist_ok=True)

    # Step 1: Export slide images
    export_img_dir = os.path.join(output_dir, "slides_img")
    os.makedirs(export_img_dir, exist_ok=True)
    export_slides_as_images(pptx_path, export_img_dir)

    # Step 2: Process each slide image (remove logo, extract heading area)
    slide_image_map = {}
    slide_heading_map = {}

    filenames = [
    f for f in os.listdir(export_img_dir)
    if re.search(r'\d+', os.path.splitext(f)[0].split("_")[-1])
    ]
    filenames.sort(key=lambda x: int(re.search(r'\d+', os.path.splitext(x)[0].split("_")[-1]).group()))


    for filename in filenames[1:-1]:
        img_path = os.path.join(export_img_dir, filename)
        slide_no = os.path.splitext(filename)[0].split("_")[-1]
        slide_no = re.search(r'\d+', slide_no).group()
        
        cropped_heading_img_path = os.path.join(output_dir, "Heading_img")
        cropped_img_path = os.path.join(output_dir, "Cropped_img")
        os.makedirs(cropped_heading_img_path, exist_ok=True)
        os.makedirs(cropped_img_path, exist_ok=True)
        
        # Create paths for cropped images
        cropped_main_path = os.path.join(cropped_img_path, f"cropped_slide_{slide_no}.png")
        heading_img_path = os.path.join(cropped_heading_img_path, f"heading_slide_{slide_no}.png")
        
        #print(f"Processing {img_path}...")
        main_img_path, heading_img_path = remove_logo_and_extract_heading(
            img_path,
            save_main_path=cropped_main_path,
            save_heading_path=heading_img_path
        )

        slide_key = f"slide_{int(slide_no)}"
        slide_image_map[slide_key] = main_img_path
        slide_heading_map[slide_key] = heading_img_path

    # Step 3: Extract text from slides
    all_text = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_text.append(text)

        combined_text = "\n".join(slide_text)
        all_text.append(f"Slide {i+1}:\n{combined_text}")

    full_text = "\n\n".join(all_text)
    
    # Step 4: Extract headings from images
    slide_headings_text = {}
    
    for slide_key, heading_img_path in slide_heading_map.items():
        img = Image.open(heading_img_path)
        text = pytesseract.image_to_string(img).strip()
        slide_headings_text[slide_key] = text

    # Step 4: Extract metadata with GPT
    

    # best_img = slide_image_map.get(info["best_image_slide"])
    # if best_img and os.path.exists(best_img):
    #     doc.add_paragraph("Machine Overview:")
    #     doc.add_picture(best_img, width=Inches(5))

    # # Step 7: Add slide content with headings and images
    # doc.add_page_break()
    #

def process_dap_to_docx(folder_path, template_path):
    """
    Extracts headings and images from a DAP PowerPoint and inserts them into a DOCX template.

    Args:
        folder_path (str): Path to the folder containing the DAP PowerPoint file.
        template_path (str): Path to the base DOCX template.
        output_docx_path (str): Path to save the final generated manual.
        placeholder (str): Placeholder text in the DOCX to be replaced.
    """
    pptx_files = glob.glob(os.path.join(folder_path, "*.pptx"))
    if len(pptx_files) == 0:
        print("⚠️ No PPTX file found in DAP folder. Aborting.")
        return
    elif len(pptx_files) > 1:
        print("⚠️ Multiple PPTX files found in DAP folder. Aborting.")
        return

    pptx_path = pptx_files[0]
    # print(f"📂 Found PPTX file: {pptx_path}")
    
    output_dir = os.path.abspath("dap_img_extracted")
    placeholder="{{MACHINE_OVERVIEW_DAP}}"
    # Step 1: Extract slide data (images + headings via OCR)
    print("Extracting DAP slide data...")
    extract_DAP_text_and_images(pptx_path, output_dir)
    
    print("Inserting DAP slide data into DOCX...")
    # Step 2: Insert extracted data into DOCX
    insert_slide_content_at_placeholder(
        template_path=template_path,
        output_path=template_path,
        placeholder=placeholder,
        headings_dict=slide_headings_text,
        image_dict=slide_image_map
    )

def process_sop_to_docx(folder_path, template_path):
    """
    Extracts headings and images from a SOP PowerPoint and inserts them into a DOCX template.

    Args:
        pptx_path (str): Path to the SOP PowerPoint file.
        template_path (str): Path to the base DOCX template.
        output_docx_path (str): Path to save the final generated manual.
        placeholder (str): Placeholder text in the DOCX to be replaced.
    """
    pptx_files = glob.glob(os.path.join(folder_path, "*.pptx"))
    if len(pptx_files) == 0:
        print("⚠️ No PPTX file found in SOP folder. Aborting.")
        return
    elif len(pptx_files) > 1:
        print("⚠️ Multiple PPTX files found in SOP folder. Aborting.")
        return

    pptx_path = pptx_files[0]
    # print(f"📂 Found PPTX file: {pptx_path}")
    
    output_dir = os.path.abspath("sop_img_extracted")
    placeholder="{{Upload_SOP_here}}"
    
    # Step 1: Extract slide data (images + headings via OCR)
    print("Extracting SOP slide data...")
    extract_DAP_text_and_images(pptx_path, output_dir)
    
    print("Inserting SOP slide data into DOCX...")
    # Step 2: Insert extracted data into DOCX
    insert_slide_content_at_placeholder(
        template_path=template_path,
        output_path=template_path,
        placeholder=placeholder,
        headings_dict=slide_headings_text,
        image_dict=slide_image_map
    )

def process_hmi_to_docx(folder_path, template_path):
    """
    Extracts headings and images from a HMI PowerPoint and inserts them into a DOCX template.

    Args:
        folder_path (str): Path to the folder containing the HMI PowerPoint file.
        template_path (str): Path to the base DOCX template.
        output_docx_path (str): Path to save the final generated manual.
        placeholder (str): Placeholder text in the DOCX to be replaced.
    """
    pptx_files = glob.glob(os.path.join(folder_path, "*.pptx"))
    if len(pptx_files) == 0:
        print("⚠️ No PPTX file found in HMI folder. Aborting.")
        return
    elif len(pptx_files) > 1:
        print("⚠️ Multiple PPTX files found in HMI folder. Aborting.")
        return

    pptx_path = pptx_files[0]
    # print(f"📂 Found PPTX file: {pptx_path}")
    
    output_dir = os.path.abspath("hmi_img_extracted")
    placeholder="{{Upload_HMI_here}}"

    # Step 1: Extract slide data (images + headings via OCR)
    print("Extracting HMI slide data...")
    extract_DAP_text_and_images(pptx_path, output_dir)
    
    print("Inserting HMI slide data into DOCX...")
    # Step 2: Insert extracted data into DOCX
    insert_slide_content_at_placeholder(
        template_path=template_path,
        output_path=template_path,
        placeholder=placeholder,
        headings_dict=slide_headings_text,
        image_dict=slide_image_map
    )

def process_scada_to_docx(folder_path, template_path):
    """
    Extracts headings and images from a SCADA PowerPoint and inserts them into a DOCX template.

    Args:
        folder_path (str): Path to the folder containing the SCADA PowerPoint file.
        template_path (str): Path to the base DOCX template.
        output_docx_path (str): Path to save the final generated manual.
        placeholder (str): Placeholder text in the DOCX to be replaced.
    """
    pptx_files = glob.glob(os.path.join(folder_path, "*.pptx"))
    if len(pptx_files) == 0:
        print("⚠️ No PPTX file found in SCADA folder. Aborting.")
        return
    elif len(pptx_files) > 1:
        print("⚠️ Multiple PPTX files found in SCADA folder. Aborting.")
        return

    pptx_path = pptx_files[0]
    # print(f"📂 Found PPTX file: {pptx_path}")
    
    output_dir = os.path.abspath("scada_img_extracted")
    placeholder="{{Upload_scada_screens_here}}"

    # Step 1: Extract slide data (images + headings via OCR)
    print("Extracting SCADA slide data...")
    extract_DAP_text_and_images(pptx_path, output_dir)
    
    print("Inserting SCADA slide data into DOCX...")
    # Step 2: Insert extracted data into DOCX
    insert_slide_content_at_placeholder(
        template_path=template_path,
        output_path=template_path,
        placeholder=placeholder,
        headings_dict=slide_headings_text,
        image_dict=slide_image_map
    )




